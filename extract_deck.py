import json
import collections
import collections.abc
from pptx import Presentation
import shutil

def extract_deck_to_json(pptx_path, json_path):
    prs = Presentation(pptx_path)
    deck_data = []

    for idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": idx + 1,
            "title": "",
            "content": []
        }
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Usually the first shape is the title, or we can check shape type
            text = shape.text.strip()
            if not text:
                continue

            if shape == slide.shapes[0] and not slide_data["title"]:
                slide_data["title"] = text
            else:
                # Add bullets or paragraphs
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text.strip())
                if paragraphs:
                    slide_data["content"].extend(paragraphs)

        deck_data.append(slide_data)
    
    with open(json_path, 'w') as f:
        json.dump(deck_data, f, indent=2)
    print(f"Extracted {len(deck_data)} slides to {json_path}")
    
    # Also copy the raw pptx for the native viewer
    public_pptx = "dashboard/public/exec_deck.pptx"
    shutil.copy2(pptx_path, public_pptx)
    print(f"Copied {pptx_path} to {public_pptx}")

if __name__ == "__main__":
    import sys
    pptx_path = "output/exec_deck.pptx"
    json_path = "dashboard/public/deck.json"
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    if len(sys.argv) > 2:
        json_path = sys.argv[2]
    extract_deck_to_json(pptx_path, json_path)
