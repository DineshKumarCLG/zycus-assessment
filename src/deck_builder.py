"""
Deck Builder — converts SynthesisResult into a real .pptx executive deck.

Uses python-pptx to produce a highly styled, modern executive presentation with 
structured layouts, custom typography, and visual dashboard-style cards.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.models import SynthesisResult

logger = logging.getLogger(__name__)

# Premium Brand Color Palette
COLOR_BG_DARK = RGBColor(0x0F, 0x17, 0x2A)     # Deep Slate 900
COLOR_BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)    # Slate 50 (Slide Background)
COLOR_TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)   # Slate 800 (Header Text)
COLOR_TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)  # Slate 500
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ACCENT = RGBColor(0x4F, 0x46, 0xE5)     # Indigo 600

# RAG Semantic Theme Colors (BG, Border, Text)
THEME_RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
THEME_RED_BORDER = RGBColor(0xFE, 0xCA, 0xCA)
THEME_RED_TEXT = RGBColor(0x99, 0x1B, 0x1B)

THEME_AMBER_BG = RGBColor(0xFF, 0xF7, 0xED)
THEME_AMBER_BORDER = RGBColor(0xFE, 0xD7, 0xAA)
THEME_AMBER_TEXT = RGBColor(0x9A, 0x34, 0x12)

THEME_GREEN_BG = RGBColor(0xF0, 0xFD, 0xFA)
THEME_GREEN_BORDER = RGBColor(0x99, 0xF6, 0xE4)
THEME_GREEN_TEXT = RGBColor(0x06, 0x5F, 0x46)

THEME_NEUTRAL_BG = RGBColor(0xFF, 0xFF, 0xFF)   # Card White
THEME_NEUTRAL_BORDER = RGBColor(0xE2, 0xE8, 0xF0) # Slate 200
THEME_NEUTRAL_TEXT = RGBColor(0x33, 0x41, 0x55)  # Slate 700

# Typography Settings
FONT_TITLE = "Georgia"
FONT_BODY = "Segoe UI"


def _add_title_slide(prs: Any, title: str, subtitle_bullets: list[str]) -> None:
    """Add a dark, high-end title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout

    # Apply solid dark slate background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()

    # Left decorative accent block
    accent_bar = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(0.08), Inches(2.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()

    # Title text box (separate from subtitle)
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(8.2), Inches(2.0))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_right = 0

    # Title paragraph
    p_title = title_tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

    # Subtitle text box (placed below the title, with extra height to prevent clipping)
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(8.2), Inches(1.4))
    sub_tf = sub_box.text_frame
    sub_tf.word_wrap = True
    sub_tf.margin_left = sub_tf.margin_top = sub_tf.margin_bottom = sub_tf.margin_right = 0

    # Subtitle paragraphs
    for i, bullet in enumerate(subtitle_bullets):
        p_sub = sub_tf.paragraphs[0] if i == 0 else sub_tf.add_paragraph()
        p_sub.text = bullet
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p_sub.space_before = Pt(4)


def _format_card_text(paragraph: Any, text: str, text_color: RGBColor) -> None:
    """Formats bullet text inside a card to bold headers and create visual hierarchy."""
    paragraph.font.name = FONT_BODY
    paragraph.font.size = Pt(14)
    paragraph.line_spacing = 1.15
    
    # Check if there is a header segment (split by colon ':')
    if ":" in text:
        header, details = text.split(":", 1)
        # Bold Header Run
        r1 = paragraph.add_run()
        r1.text = header + ":"
        r1.font.bold = True
        r1.font.color.rgb = text_color
        # Regular Details Run
        r2 = paragraph.add_run()
        r2.text = details
        r2.font.bold = False
        r2.font.color.rgb = text_color
    else:
        # Fallback: bold the first 4 words if no colon
        words = text.split(" ")
        if len(words) > 4:
            header = " ".join(words[:4])
            details = " " + " ".join(words[4:])
            r1 = paragraph.add_run()
            r1.text = header
            r1.font.bold = True
            r1.font.color.rgb = text_color
            
            r2 = paragraph.add_run()
            r2.text = details
            r2.font.bold = False
            r2.font.color.rgb = text_color
        else:
            r = paragraph.add_run()
            r.text = text
            r.font.bold = False
            r.font.color.rgb = text_color


def _add_content_slide(
    prs: Any,
    title: str,
    bullets: list[str],
    notes: str | None = None,
) -> None:
    """Add a content slide styled as a modern dashboard using card components."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout

    # Apply solid light slide background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_LIGHT
    bg.line.fill.background()

    # 1. Title text box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_right = 0
    
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.name = FONT_TITLE
    title_p.font.size = Pt(26)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TEXT_DARK

    # Subtitle decorative accent line
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.15), Inches(8.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEXT_MUTED
    line.line.fill.background()

    # 2. Dynamic card generation for bullets
    N = len(bullets)
    if N > 0:
        start_top = 1.4
        max_height = 3.6
        gap = 0.12
        card_height = (max_height - (gap * (N - 1))) / N
        card_height = min(card_height, 1.4) # Limit maximum height for short lists

        for idx, bullet_text in enumerate(bullets):
            top_pos = start_top + idx * (card_height + gap)
            
            # Semantic RAG styling detection
            lower_text = bullet_text.lower()
            if "red" in lower_text or "distress" in lower_text or "critical" in lower_text or "⚠" in lower_text:
                bg_color = THEME_RED_BG
                border_color = THEME_RED_BORDER
                text_color = THEME_RED_TEXT
            elif "amber" in lower_text or "yellow" in lower_text or "risk" in lower_text or "vulnerability" in lower_text:
                bg_color = THEME_AMBER_BG
                border_color = THEME_AMBER_BORDER
                text_color = THEME_AMBER_TEXT
            elif "green" in lower_text or "stable" in lower_text or "reconciled" in lower_text:
                bg_color = THEME_GREEN_BG
                border_color = THEME_GREEN_BORDER
                text_color = THEME_GREEN_TEXT
            else:
                bg_color = THEME_NEUTRAL_BG
                border_color = THEME_NEUTRAL_BORDER
                text_color = THEME_NEUTRAL_TEXT

            # Create card background rectangle
            card = slide.shapes.add_shape(
                1, # Rectangle
                Inches(0.8), Inches(top_pos), Inches(8.4), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = bg_color
            card.line.color.rgb = border_color
            card.line.width = Pt(1)

            # Insert styled text frame in card
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.12)
            tf.margin_bottom = Inches(0.12)

            p = tf.paragraphs[0]
            _format_card_text(p, bullet_text, text_color)

    # 3. Add presenter speaker notes
    if notes and slide.has_notes_slide:
        slide.notes_slide.notes_text_frame.text = notes


def build_deck(
    synthesis: SynthesisResult,
    output_path: Path = Path("output/exec_deck.pptx"),
) -> Path:
    """Generate a high-end .pptx executive deck from the synthesis output."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    # Set slide size to widescreen 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    if not synthesis.slides:
        logger.warning("No slide content in synthesis — generating minimal deck")
        _add_title_slide(prs, "Project Health Report", ["No data available"])
        prs.save(str(output_path))
        return output_path

    # Slide 1: Title
    first = synthesis.slides[0]
    _add_title_slide(prs, first.title, first.bullets)

    # Subsequent slides: Dashboards
    for slide_content in synthesis.slides[1:]:
        _add_content_slide(
            prs,
            slide_content.title,
            slide_content.bullets,
            notes=slide_content.notes
        )

    slide_count = len(prs.slides)
    logger.info("Generated deck with %d slides", slide_count)

    prs.save(str(output_path))
    logger.info("Deck saved to %s", output_path)
    return output_path
